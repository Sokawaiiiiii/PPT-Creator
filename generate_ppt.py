import os
from openai import OpenAI
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def generate_slide_content(topic: str, slides: int = 8) -> str:
    """Generates slide content using OpenAI API."""
    prompt = f"""
    Create a structured PowerPoint outline with {slides} slides on the topic "{topic}".
    Each slide should include:
    - A clear slide title
    - 3 to 5 concise bullet points explaining the concept
    Format:
    Slide 1: Title
    - Point 1
    - Point 2
    - Point 3
    Slide 2: Title
    ...
    """
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response.choices[0].message.content.strip()

def create_ppt(content: str, filename: str = "output.pptx") -> str:
    """Converts AI-generated text into a PowerPoint file."""
    prs = Presentation()
    slides = content.split("Slide ")
    for s in slides[1:]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = [line.strip() for line in s.split("\n") if line.strip()]
        title = lines[0]
        body = "\n".join(lines[1:])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    prs.save(filename)
    return filename
